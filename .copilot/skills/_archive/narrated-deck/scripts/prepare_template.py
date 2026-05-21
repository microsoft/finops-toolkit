#!/usr/bin/env python3
"""Strip all slides from a PPTX template, leaving only masters and layouts.

This is run once to prepare a clean source-template for the skill to use.
The release-deck/source-template.pptx ships with 44 leftover content slides
from prior V8 builds; those slides' XML parts collide with new slides added
at runtime, causing PPTX corruption ("PowerPoint found a problem with content
in the file" prompt).

Usage:
    python3 prepare_template.py <input.pptx> <output.pptx>

The output template retains:
- All slide masters (theme, fonts, default colors)
- All slide layouts (so build_pptx can pick TITLE / BULLETS / TABLE / etc.)
- Theme assets (logo, brand graphics) attached to masters

The output template removes:
- Every concrete slide
- The slide parts they reference (both the XML in the package AND the rels)

Implementation: python-pptx's `drop_rel` removes the relationship but leaves
orphaned XML parts in the OPC package (the .pptx zip). To truly clean the
template, we post-process the saved file with `zipfile`, removing any
ppt/slides/slide*.xml and ppt/slides/_rels/slide*.xml.rels entries that
aren't referenced from the presentation's main rels.
"""
import re
import sys
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except ImportError:
    print("ERROR: install python-pptx (pip install python-pptx)", file=sys.stderr)
    sys.exit(2)


def strip_slides_python_pptx(input_path: Path, intermediate_path: Path) -> int:
    """Step 1: drop slide rels via python-pptx. Returns count of slides removed."""
    prs = Presentation(str(input_path))
    slides_before = len(prs.slides)
    sldIdLst = prs.slides._sldIdLst
    rIds = []
    for sldId in list(sldIdLst):
        rIds.append(sldId.get(qn("r:id")))
        sldIdLst.remove(sldId)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    intermediate_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(intermediate_path))
    return slides_before


def strip_orphan_parts(intermediate_path: Path, output_path: Path) -> int:
    """Step 2: rewrite the .pptx zip, dropping orphaned slide parts.

    Reads the package, identifies any ppt/slides/* entry not referenced from
    ppt/_rels/presentation.xml.rels, and skips them when writing the new zip.
    Returns count of orphan files removed.
    """
    with zipfile.ZipFile(intermediate_path, "r") as zin:
        names = zin.namelist()
        try:
            pres_rels = zin.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        except KeyError:
            pres_rels = ""
        try:
            content_types = zin.read("[Content_Types].xml").decode("utf-8")
        except KeyError:
            content_types = ""

        def is_referenced(slide_path: str) -> bool:
            slide_filename = slide_path.split("/")[-1]
            return (
                slide_filename in pres_rels
                or f"slides/{slide_filename}" in pres_rels
                or f"/ppt/slides/{slide_filename}" in content_types
            )

        slide_pattern = re.compile(r"^ppt/slides/slide\d+\.xml$")
        slide_rel_pattern = re.compile(r"^ppt/slides/_rels/slide\d+\.xml\.rels$")
        orphan_paths = set()
        for name in names:
            if slide_pattern.match(name):
                if not is_referenced(name):
                    orphan_paths.add(name)
            elif slide_rel_pattern.match(name):
                slide_companion = name.replace("/_rels/", "/").replace(".rels", "")
                if not is_referenced(slide_companion):
                    orphan_paths.add(name)

        kept_content_types = content_types
        for orphan in orphan_paths:
            if not orphan.endswith(".rels"):
                pattern = re.compile(
                    rf'<Override\s+PartName="/{re.escape(orphan)}"[^/]*/>'
                )
                kept_content_types = pattern.sub("", kept_content_types)

        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                if name in orphan_paths:
                    continue
                if name == "[Content_Types].xml":
                    zout.writestr(name, kept_content_types)
                else:
                    zout.writestr(name, zin.read(name))
    return len(orphan_paths)


def main():
    if len(sys.argv) != 3:
        print("Usage: prepare_template.py <input.pptx> <output.pptx>", file=sys.stderr)
        sys.exit(2)
    inp = Path(sys.argv[1])
    out = Path(sys.argv[2])
    if not inp.exists():
        print(f"Input not found: {inp}", file=sys.stderr)
        sys.exit(2)
    intermediate = out.with_suffix(".intermediate.pptx")
    slides_removed = strip_slides_python_pptx(inp, intermediate)
    orphans_removed = strip_orphan_parts(intermediate, out)
    intermediate.unlink()
    after = Presentation(str(out))
    in_size = inp.stat().st_size
    out_size = out.stat().st_size
    print(
        f"Stripped {slides_removed} slides via python-pptx, "
        f"removed {orphans_removed} orphan files from zip"
    )
    print(
        f"Result: {len(after.slides)} slides, {len(after.slide_layouts)} layouts, "
        f"{len(after.slide_masters)} masters"
    )
    print(f"Size: {in_size:,} → {out_size:,} bytes ({out_size / in_size * 100:.0f}%)")


if __name__ == "__main__":
    main()
