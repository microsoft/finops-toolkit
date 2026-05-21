#!/usr/bin/env python3
"""Build a branded FinOps Toolkit PowerPoint deck from a scene plan.

Reads either:
  - scenes.json directly (default — fully automated path)
  - A deck-outline.md table (human-edited path — same format as release-deck)

Renders one slide per scene using simple layouts:
  TITLE    — hero slide with H1 + subtitle
  BULLETS  — heading + bullet list
  TABLE    — heading + a real PPTX table from the source markdown table
  CODE     — heading + monospace text from fenced code block
  CALLOUT  — heading + large pull quote
  OUTRO    — closing slide

Brand identity comes from the bundled assets/source-template.pptx, which carries
the FinOps Toolkit purple gradient (#7B33C0 → #4E2E7A), the hexagon logo, and
Microsoft Learn-style typography. Slides inherit this template's masters.

Speaker notes are populated from the scene's normalized narration, so the .pptx
is a complete, self-contained artifact: open in PowerPoint and you'll see the
slide and the matching narration.

This is a generalized version of src/templates/sre-agent/training/release-deck/build.py
with the same template + python-pptx mechanics, but driven by any docs-mslearn doc
rather than a hand-tuned outline.
"""
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print(
        "ERROR: python-pptx not installed. Install with: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


SCRIPT_DIR = Path(__file__).parent
sys.path.insert(0, str(SCRIPT_DIR))  # so we can import normalize_text
ASSETS_DIR = SCRIPT_DIR.parent / "assets"
TEMPLATE_PATH = ASSETS_DIR / "source-template-clean.pptx"

# Brand colors (FinOps Toolkit)
PURPLE_DARK = RGBColor(0x4E, 0x2E, 0x7A)
PURPLE_BRIGHT = RGBColor(0x7B, 0x33, 0xC0)
CHARCOAL = RGBColor(0x1F, 0x23, 0x28)
LIGHT_GRAY = RGBColor(0x80, 0x80, 0x80)


def find_layout(prs: "Presentation", names: list[str]):
    """Find a layout by trying each candidate name in order. Falls back to layout 0."""
    for name in names:
        for layout in prs.slide_layouts:
            if layout.name and layout.name.lower() == name.lower():
                return layout
    return prs.slide_layouts[0]


def find_body_placeholder(slide):
    """Return the first non-title placeholder on a slide.

    Different brand-template layouts use different idx values for the body
    placeholder (Title Slide uses idx=12, Title and Content uses idx=10).
    Looking up by idx is fragile. Looking up by "first non-title placeholder"
    is robust: the title is always type TITLE (1) or CENTER_TITLE (13), and
    everything else is body content.
    """
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type
        # 1 = TITLE, 13 = CENTER_TITLE, 15 = SUB_TITLE — skip all of these
        if ph_type in (1, 13, 15):
            continue
        return ph
    return None


def add_text_to_placeholder(placeholder, text: str, font_size: int = 18, bold: bool = False, color=None):
    """Set text into a placeholder's text frame with style."""
    if placeholder is None:
        return
    tf = placeholder.text_frame
    tf.clear()
    paras = text.split("\n")
    for i, para_text in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def add_speaker_notes(slide, notes_text: str):
    """Populate the slide's speaker-notes pane with the narration.

    The text passed in is already in spoken voice (normalized upstream by
    doc-to-yaml). We just clean any HTML break markers and write it to the
    notes pane verbatim.
    """
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    paras = re.split(r"\n\s*\n|<br\s*/>\s*<br\s*/>", notes_text)
    for i, p in enumerate(paras):
        para_text = re.sub(r"<break\s+time[^/]*/>", "", p)
        para_text = re.sub(r"<br\s*/>", " ", para_text).strip()
        if not para_text:
            continue
        para = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        run = para.add_run()
        run.text = para_text
        run.font.size = Pt(12)


def render_title_slide(prs, scene, h1):
    layout = find_layout(prs, ["Title Slide", "Title", "Section Header"])
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = h1
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = CHARCOAL
            run.font.bold = True
    # Subtitle = first sentence of the scene text
    text = scene.get("text", "").strip()
    first_sentence = re.split(r"(?<=[.!?])\s+", text, maxsplit=1)[0] if text else ""
    body_ph = find_body_placeholder(slide)
    if body_ph is not None:
        add_text_to_placeholder(body_ph, first_sentence, font_size=20, color=LIGHT_GRAY)
    add_speaker_notes(slide, scene.get("text", ""))
    return slide


def render_bullets_slide(prs, scene):
    layout = find_layout(prs, ["Title and Content", "Title and Body", "Content"])
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = scene["heading"]
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = CHARCOAL
    # Body: extract bullets from raw markdown
    bullets = extract_bullets_from_markdown(scene["raw_markdown"])
    if not bullets:
        # Fall back to first 3 sentences of the cleaned text
        sentences = re.split(r"(?<=[.!?])\s+", scene["text"].strip())
        bullets = [s for s in sentences[:5] if s]
    body_ph = find_body_placeholder(slide)
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
    add_speaker_notes(slide, scene.get("text", ""))
    return slide


def render_table_slide(prs, scene):
    layout = find_layout(prs, ["Title and Content", "Title Only", "Content"])
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = scene["heading"]
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = CHARCOAL
    rows = parse_markdown_table(scene["raw_markdown"])
    if rows:
        n_rows, n_cols = len(rows), len(rows[0])
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.4 * n_rows)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table
        for r, row in enumerate(rows):
            for c, cell_text in enumerate(row):
                cell = table.cell(r, c)
                cell.text = cell_text
                # Header row formatting
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(11 if r > 0 else 12)
                        run.font.bold = (r == 0)
                        if r == 0:
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_speaker_notes(slide, scene.get("text", ""))
    return slide


def render_code_slide(prs, scene):
    layout = find_layout(prs, ["Title and Content", "Title Only"])
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = scene["heading"]
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = CHARCOAL
    code = extract_first_code_block(scene["raw_markdown"])
    if code:
        body_ph = find_body_placeholder(slide)
        if body_ph is not None:
            tf = body_ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = code[:1500]  # cap to avoid overflow
            run.font.name = "Consolas"
            run.font.size = Pt(11)
            run.font.color.rgb = CHARCOAL
    add_speaker_notes(slide, scene.get("text", ""))
    return slide


def render_callout_slide(prs, scene):
    layout = find_layout(prs, ["Section Header", "Title Slide", "Title"])
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = scene["heading"]
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = CHARCOAL
    text = scene["text"].strip()
    body_ph = find_body_placeholder(slide)
    if body_ph is not None:
        add_text_to_placeholder(body_ph, text, font_size=22, color=PURPLE_DARK)
    add_speaker_notes(slide, text)
    return slide


def render_outro_slide(prs, h1: str, next_title: str = "", scene: dict = None):
    layout = find_layout(prs, ["Section Header", "Title Slide"])
    slide = prs.slides.add_slide(layout)
    title_text = "End of module" if not next_title else "Coming up next"
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = CHARCOAL
    subtitle = next_title if next_title else h1
    body_ph = find_body_placeholder(slide)
    if body_ph is not None:
        add_text_to_placeholder(body_ph, subtitle, font_size=24, color=PURPLE_BRIGHT, bold=True)
    if scene and scene.get("text"):
        notes_text = scene["text"]
        if next_title:
            notes_text = f"Coming up next: {next_title}."
    else:
        notes_text = (
            f"Coming up next: {next_title}." if next_title
            else f"That concludes {h1}. Find related modules on Microsoft Learn under FinOps Toolkit."
        )
    add_speaker_notes(slide, notes_text)
    return slide


def extract_bullets_from_markdown(raw: str) -> list[str]:
    """Pull bullet items from raw markdown; strip link/bold/italic/break syntax."""
    bullets = re.findall(r"^\s*-\s+(.+)$", raw, flags=re.MULTILINE)
    cleaned = []
    for b in bullets:
        b = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", b)
        b = re.sub(r"\*\*([^*]+)\*\*", r"\1", b)
        b = re.sub(r"\*([^*]+)\*", r"\1", b)
        b = re.sub(r"_([^_\n]+)_", r"\1", b)
        b = re.sub(r"<br\s*/?>", " \u2014 ", b)
        b = re.sub(r"`([^`]+)`", r"\1", b)
        b = re.sub(r"\s+", " ", b).strip()
        cleaned.append(b)
    return cleaned


def parse_markdown_table(raw: str) -> list[list[str]]:
    """Find the first markdown table in `raw` and return as rows of cells."""
    lines = raw.split("\n")
    table_lines = []
    in_table = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            in_table = True
            table_lines.append(stripped)
        elif in_table:
            break  # table ended
    if not table_lines:
        return []
    rows = []
    for line in table_lines:
        # Skip the separator row (|---|---|)
        if re.match(r"^\|[\s:|-]+\|$", line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        # Strip markdown link syntax inside cells
        cells = [re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", c) for c in cells]
        cells = [re.sub(r"`([^`]+)`", r"\1", c) for c in cells]
        rows.append(cells)
    return rows


def extract_first_code_block(raw: str) -> str:
    """Return the contents of the first fenced code block in raw markdown."""
    m = re.search(r"```[a-z]*\n(.*?)\n```", raw, flags=re.DOTALL)
    return m.group(1) if m else ""


def build(scenes_data: dict, out_path: Path, next_title: str = "") -> Path:
    if not TEMPLATE_PATH.exists():
        print(f"ERROR: missing source template at {TEMPLATE_PATH}", file=sys.stderr)
        sys.exit(2)
    prs = Presentation(str(TEMPLATE_PATH))
    h1 = scenes_data["h1"]
    for scene in scenes_data["scenes"]:
        layout = scene["layout"]
        if layout == "TITLE":
            render_title_slide(prs, scene, h1)
        elif layout == "TABLE":
            render_table_slide(prs, scene)
        elif layout == "CODE":
            render_code_slide(prs, scene)
        elif layout == "CALLOUT":
            render_callout_slide(prs, scene)
        elif layout == "OUTRO":
            render_outro_slide(prs, h1, next_title=next_title)
        else:  # BULLETS or REFERENCE — both render as bullets-with-heading
            render_bullets_slide(prs, scene)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main():
    args = sys.argv[1:]
    if "--out" not in args:
        print(
            "Usage: build_pptx.py --out <path.pptx> [--scenes <path>] [--next <title>]",
            file=sys.stderr,
        )
        sys.exit(2)
    out_path = Path(args[args.index("--out") + 1])
    if "--scenes" in args:
        scenes_data = json.loads(Path(args[args.index("--scenes") + 1]).read_text())
    else:
        scenes_data = json.loads(sys.stdin.read())
    next_title = args[args.index("--next") + 1] if "--next" in args else ""
    result = build(scenes_data, out_path, next_title=next_title)
    print(result)


if __name__ == "__main__":
    main()
