"""Build a FinOps Toolkit deck (PPTX) from a directory of per-slide YAML files.

All paths are required — no defaults, no environment-variable fallbacks.
Pass each path explicitly so the skill works in any project / any cwd.

Usage:
    python build_yaml.py \
        --slides-dir  /path/to/slides/                 # *.yaml inputs
        --template    $SKILL/assets/source-template.pptx
        --labels      $SKILL/assets/labels.yaml
        --output      /path/to/out.pptx
        [--part 1|2|3|all]                             # default 'all'

`--part` only filters slides whose `id:` matches the V8 deck cluster pattern
(P1.x for part 1, P2.x for part 2, H/Z for part 3). For non-V8 decks, leave
it 'all'.
"""
import argparse
import os
import re
import sys
import tempfile
import time
from pathlib import Path

import yaml
from pptx import Presentation


def parse_slides(slides_dir: Path, _b):
    """Read slides_dir/*.yaml in filename order, return rows in build.py schema."""
    rows = []
    files = sorted(slides_dir.glob("*.yaml"))
    if not files:
        sys.exit(f"ERROR: no YAML files found in {slides_dir}/")
    for f in files:
        s = yaml.safe_load(f.read_text(encoding="utf-8"))
        bullets = [_b.strip_md(_b.strip_bullet(b)) for b in _b.html_br_to_lines(s.get("content", "") or "")]
        notes_text = s.get("notes", "") or ""
        notes_paras = [_b.strip_md(p) for p in _b.html_br_to_lines(notes_text)] or [_b.strip_md(notes_text)]
        layout = s.get("layout", {}) or {}
        kind = layout.get("kind", "BULLETS")
        addresses = layout.get("addresses") or []
        verdict = layout.get("verdict")
        assets = {}
        screens = s.get("screens") or ""
        if screens:
            am = re.search(r"assets:\s*(.+?)$", screens)
            if am:
                for kv in re.finditer(r"(\w+)\s*=\s*([^\s|]+)", am.group(1)):
                    assets[kv.group(1)] = kv.group(2)
        rows.append({
            "num": s["id"],
            "stage": _b.strip_md(s.get("cluster", "") or ""),
            "title": _b.strip_md(s.get("title", "") or ""),
            "bullets": bullets,
            "notes": notes_paras,
            "kind": kind,
            "layout": "",
            "assets": assets,
            "addresses": addresses,
            "verdict": verdict.lower() if verdict else None,
        })
    return rows


def main():
    parser = argparse.ArgumentParser(description="Build a deck from per-slide YAML.")
    parser.add_argument("--slides-dir", required=True, help="Directory containing per-slide *.yaml files")
    parser.add_argument("--template", required=True, help="Path to source-template.pptx")
    parser.add_argument("--labels", required=True, help="Path to labels.yaml (branding strings)")
    parser.add_argument("--output", required=True, help="Path to write output .pptx")
    parser.add_argument("--charts", help="Optional: path to charts dir (required if slides use screens: chart=...)")
    parser.add_argument("--assets", help="Optional: path to image-assets dir (required if slides use screens: image=...)")
    parser.add_argument(
        "--part",
        choices=["1", "2", "3", "all"],
        default="all",
        help="V8-deck filter: 'all' (default), '1' (P1.x), '2' (P2.x), '3' (H/Z). Leave 'all' for non-V8 decks.",
    )
    args = parser.parse_args()

    slides_dir = Path(args.slides_dir).expanduser().resolve()
    if not slides_dir.is_dir():
        sys.exit(f"ERROR: --slides-dir not a directory: {slides_dir}")
    out_path = Path(args.output).expanduser().resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # build.py must be configured before it's safe to use anything from it.
    # Add this script's dir to sys.path so the import resolves to our copy.
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    import build as _b
    _b.configure(
        template=args.template,
        labels=args.labels,
        charts_dir=args.charts,
        assets_dir=args.assets,
        output=str(out_path),
    )

    rows = parse_slides(slides_dir, _b)

    if args.part != "all":
        pattern, _ = _b.PART_FILTERS[args.part]
        prog = re.compile(pattern) if isinstance(pattern, str) else pattern
        before = len(rows)
        rows = [r for r in rows if prog.match(r["num"])]
        print(f"Filtered to part {args.part}: {len(rows)} of {before} slides match {prog.pattern}")

    if not rows:
        sys.exit(f"ERROR: no rows matched filter for part {args.part}")

    prs = Presentation(str(_b.SOURCE))

    for kind in {r["kind"] for r in rows}:
        if kind not in _b.KIND_LAYOUT:
            sys.exit(f"ERROR: layout kind '{kind}' not in build.py KIND_LAYOUT. Valid: {sorted(_b.KIND_LAYOUT)}")
        _b.find_layout(prs, _b.KIND_LAYOUT[kind])

    _b.delete_all_slides(prs)
    print("Cleared all original slides.\n")

    for row in rows:
        kind = row["kind"]
        layout = _b.find_layout(prs, _b.KIND_LAYOUT[kind])
        slide = prs.slides.add_slide(layout)
        _b.RENDERERS[kind](slide, row)
        if row["notes"] and any(n.strip() for n in row["notes"]):
            tf = slide.notes_slide.notes_text_frame
            tf.clear()
            tf.paragraphs[0].text = row["notes"][0]
            for note_para in row["notes"][1:]:
                p = tf.add_paragraph()
                p.text = note_para
        print(f"  {row['num']:5s}  [{kind:18s}]  {layout.name:40s}  {row['title'][:55]}")

    prs.core_properties.author = _b.LABELS["pptx.author"]
    prs.core_properties.last_modified_by = _b.LABELS["pptx.author"]

    out_dir = out_path.parent
    out_name_final = out_path.name
    lock = out_dir / f"~${out_name_final}"
    if lock.exists():
        alt = out_path.stem + f".rebuild-{int(time.time())}.pptx"
        out_path = out_dir / alt
        print(f"\n⚠️  PowerPoint lock detected on {out_name_final}. Writing to {alt} instead.")

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx", dir=str(out_dir))
    os.close(tmp_fd)
    try:
        prs.save(tmp_path)
        os.replace(tmp_path, str(out_path))
        print(f"\nWrote {out_path}")
        print(f"Slides: {len(rows)}")
    except Exception:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise


if __name__ == "__main__":
    main()
