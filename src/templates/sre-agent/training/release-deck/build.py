#!/usr/bin/env python3
"""Build V8 deck via python-pptx end-to-end (no XML hand-crafting).

python-pptx uses the official OPC API and produces files PowerPoint accepts
without prompting for repair.

Pipeline:
  1. Load source.pptx (template) into a Presentation
  2. Resolve the 4 layouts I need by name
  3. Delete every existing slide
  4. For each V6 row, add_slide(layout) and populate placeholders + notes
  5. Save
"""
import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree

ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "source-template.pptx"
ASSETS = ROOT / "assets"
OUT = ROOT / "finops-toolkit-sre-agent-release-training.pptx"
V8 = ROOT / "deck-outline-v8.md"
CHARTS = ROOT / "charts" / "svg"

PART_FILTERS = {
    "1": (
        re.compile(r"^(0\.1|0\.2|P1\.\d+\.[ABC]|1\.99)$"),
        "finops-toolkit-sre-agent-release-training-part1-deal-motion.pptx",
    ),
    "2": (
        re.compile(r"^(2\.0\.1|2\.0\.2|P2\.\d+\.[ABC]|2\.99)$"),
        "finops-toolkit-sre-agent-release-training-part2-operate-motion.pptx",
    ),
    "3": (
        re.compile(r"^([HZ]\.\d+)$"),
        "finops-toolkit-sre-agent-release-training-part3-honest-and-close.pptx",
    ),
}

# ─── Layout names from the template ──────────────────────────
LAYOUT_TITLE   = "Title Slide"
LAYOUT_SECTION = "Section Title"
LAYOUT_CONTENT = "Title and Content"
LAYOUT_TWOCOL  = "Two Column Content with Subheads"
LAYOUT_QUOTE   = "Quote 1"
LAYOUT_CODE    = "Developer Code Layout"

# Map kind → layout name
KIND_LAYOUT = {
    "TITLE":            LAYOUT_TITLE,
    "BULLETS":          LAYOUT_CONTENT,
    "ASK_A":            LAYOUT_CONTENT,
    "ASK_B":            LAYOUT_CONTENT,
    "ASK_C":            LAYOUT_CONTENT,
    "INDEX":            LAYOUT_CONTENT,
}

# Verdict colors for footer chip
VERDICT_COLORS = {
    "green":  ("2E7D32", "E8F5E9"),   # dark green text, light green bg
    "yellow": ("F57F17", "FFF8E1"),   # dark amber text, light amber bg
    "red":    ("C62828", "FFEBEE"),   # dark red text, light red bg
}


# ─── V8 parser (handles 3 concatenated tables) ──────────────

def parse_table_cells(line):
    if not line.startswith("|") or not line.endswith("|"):
        return None
    # Handle escaped pipes (\|) within cells by temporarily replacing them
    line_clean = line.replace("\\|", "\x00PIPE\x00")
    cells = [c.strip().replace("\x00PIPE\x00", "|") for c in line_clean[1:-1].split("|")]
    return cells


def html_br_to_lines(s):
    parts = re.split(r"<br\s*/?>", s)
    return [p.strip() for p in parts if p.strip()]


def strip_md(s):
    s = re.sub(r"\[\[([^\]|]+)(?:\|[^\]]+)?\]\]", r"\1", s)
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    # Italic: only match *text* (no spaces around) — preserves cron asterisks like "30 6 * * *"
    s = re.sub(r"(?<![\*\s])\*(\S[^*]*?\S)\*(?![\*])", r"\1", s)
    s = re.sub(r"(?<![\*\s])\*(\S)\*(?![\*])", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    s = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", s)
    # Decode HTML entities
    s = s.replace("&gt;", ">").replace("&lt;", "<").replace("&amp;", "&")
    s = s.replace("&quot;", '"').replace("&#39;", "'")
    return s


def strip_bullet(s):
    return re.sub(r"^[•·\-\*]\s+", "", s.strip())


def parse_labels():
    text = V8.read_text(encoding="utf-8")
    m = re.search(r"<!--\s*DECK-LABELS\s*\n(.*?)\n\s*-->", text, re.DOTALL)
    if not m:
        sys.exit("DECK-LABELS block not found in deck-outline-v8.md")
    labels = {}
    for line in m.group(1).strip().split("\n"):
        if ":" not in line:
            continue
        k, v = line.split(":", 1)
        labels[k.strip()] = v.strip().strip("'\"")
    return labels


LABELS = parse_labels()


def parse_v8():
    text = V8.read_text(encoding="utf-8")
    lines = text.split("\n")

    # Find ALL table header lines (V8 has 3 tables)
    table_starts = []
    for i, ln in enumerate(lines):
        if ln.startswith("| # | Cluster | Asks (verbatim) |") or ln.startswith("| # | Cluster | Title |") or ln.startswith("| # |Cluster|"):
            table_starts.append(i)

    if not table_starts:
        sys.exit("V8 slide-scaffold table not found")

    rows = []
    for table_start in table_starts:
        i = table_start + 2  # skip header + separator
        while i < len(lines):
            ln = lines[i]
            if not ln.startswith("|"):
                break
            cells = parse_table_cells(ln)
            if not cells or len(cells) < 7:
                i += 1
                continue
            # Schema 8-column (with asks-verbatim) or legacy 7-column
            if len(cells) >= 8:
                num, stage, _asks_verbatim, title, content, notes, layout, screens = cells[:8]
            else:
                num, stage, title, content, notes, layout, screens = cells[:7]
            # Strip markdown from num field (drafter output may use **P1.2.A** bold)
            num = strip_md(num).strip()
            # Accept slide IDs like 0.1, P1.1.A, P2.3.B, H.1, Z.1, 1.99, 2.99, 2.0.1
            if not re.match(r"^[\d.]+$|^P\d+\.\d+\.[ABC]$|^[HZ]\.\d+$", num) or "APPENDIX" in (stage or ""):
                i += 1
                continue
            bullets = [strip_md(strip_bullet(b)) for b in html_br_to_lines(content)]
            notes_paras = [strip_md(p) for p in html_br_to_lines(notes)] or [strip_md(notes)]
            m = re.match(r"kind=(\w+)\s*·?\s*(.*)", layout)
            kind = m.group(1) if m else "BULLETS"
            layout_desc = m.group(2) if m else strip_md(layout)

            # Parse addresses and verdict from layout directive
            addresses = []
            verdict = None
            addr_m = re.search(r"addresses=([\d,]+)", layout_desc)
            if addr_m:
                addresses = [int(x.strip()) for x in addr_m.group(1).split(",") if x.strip()]
            verd_m = re.search(r"verdict=(\w+)", layout_desc)
            if verd_m:
                verdict = verd_m.group(1).lower()

            # Parse assets directive in screens column
            screens_text = strip_md(screens)
            assets = {}
            am = re.search(r"assets:\s*(.+?)$", screens_text)
            if am:
                for kv in re.finditer(r"(\w+)\s*=\s*([^\s|]+)", am.group(1)):
                    assets[kv.group(1)] = kv.group(2)

            rows.append({
                "num": num, "stage": strip_md(stage), "title": strip_md(title),
                "bullets": bullets, "notes": notes_paras, "kind": kind,
                "layout": strip_md(layout_desc),
                "assets": assets,
                "addresses": addresses,
                "verdict": verdict,
            })
            i += 1
    return rows


# ─── Layout / placeholder helpers ────────────────────────────

def find_layout(prs, name):
    """Find a layout by name, preferring the first slide master (master 0)."""
    # First try master 0 (the branded master used by V6)
    master0 = list(prs.slide_masters)[0]
    for layout in master0.slide_layouts:
        if layout.name == name:
            return layout
    # Fallback: search all masters
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise KeyError(f"Layout {name!r} not found")


def get_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_text_frame(ph, paragraphs):
    """Populate a text-frame placeholder.

    paragraphs is a list of either strings or dicts with {text, size, bold, color}.
    """
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    for i, p in enumerate(paragraphs):
        if isinstance(p, str):
            text, attrs = p, {}
        else:
            text, attrs = p["text"], p

        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # Clear any inherited content
        for r in list(para.runs):
            r.text = ""

        run = para.add_run()
        run.text = text or ""

        font = run.font
        if attrs.get("size"):
            font.size = Pt(attrs["size"])
        if attrs.get("bold"):
            font.bold = True
        if attrs.get("color"):
            font.color.rgb = RGBColor.from_string(attrs["color"])


# ─── Content classifier helpers ──────────────────────────────

def split_label(bullet):
    """Split 'Label: rest' into (label, body). Label can contain parens, hyphens, spaces."""
    m = re.match(r"^([A-Z][\w\-\s\(\)\/]{0,40}?):\s+(.+)$", bullet)
    if m:
        return m.group(1).strip(), m.group(2)
    return None, bullet


def find_prefixed(bullets, prefixes):
    """Match bullets whose label STARTS WITH any of `prefixes` (case-insensitive).

    Examples that match prefix 'Capacity':
      'Capacity: ...'
      'Capacity output: ...'
      'Capacity (cadence): ...'
    """
    plow = [p.lower() for p in prefixes]
    for b in bullets:
        label, body = split_label(b)
        if label and any(label.lower().startswith(p) for p in plow):
            return body
    return None


def add_image_to_placeholder(slide, ph, asset_filename, sl_w_emu, sl_h_emu):
    """Insert an image inside the bounds of a placeholder, then clear the placeholder text.

    SVG assets are rasterized to PNG via rsvg-convert before insertion.
    Searches both ASSETS and CHARTS directories.
    """
    asset_path = ASSETS / asset_filename
    if not asset_path.exists():
        asset_path = CHARTS / asset_filename
    if not asset_path.exists():
        print(f"   ! asset missing: {asset_filename}")
        return

    if asset_path.suffix.lower() == ".svg":
        png_path = asset_path.with_suffix(".rasterized.png")
        if not png_path.exists():
            import subprocess
            subprocess.run(
                ["rsvg-convert", "-w", "1600", "-o", str(png_path), str(asset_path)],
                check=True,
            )
        asset_path = png_path

    left, top = ph.left, ph.top
    width, height = ph.width, ph.height
    slide.shapes.add_picture(str(asset_path), left, top, width=width, height=height)
    if ph.has_text_frame:
        ph.text_frame.clear()


# ─── Per-kind renderers (use python-pptx, not raw XML) ───────

def render_title(slide, row):
    from pptx.enum.shapes import MSO_SHAPE
    if slide.shapes.title is not None:
        title_ph = slide.shapes.title
        title_ph.text = row["title"]
        for p in title_ph.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    sub = find_prefixed(row["bullets"], [LABELS["cover.subtitle-prefix"]]) or ""
    aud = find_prefixed(row["bullets"], [LABELS["cover.audience-prefix"]]) or ""
    cob = find_prefixed(row["bullets"], [LABELS["cover.cobrand-prefix"]]) or ""
    body = get_placeholder(slide, 12)
    set_text_frame(body, [
        {"text": sub, "size": 28},
        {"text": ""},
        {"text": aud, "size": 16, "color": "595959"},
    ])
    if body is not None:
        for p in body.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    # Co-brand strip across the bottom — split label "X + Y" into two anchors
    if cob:
        sl_w = slide.part.package.presentation_part.presentation.slide_width
        sl_h = slide.part.package.presentation_part.presentation.slide_height
        # Background strip
        strip_h = Inches(0.85)
        strip_top = sl_h - strip_h
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, strip_top, sl_w, strip_h)
        strip.fill.solid()
        strip.fill.fore_color.rgb = RGBColor.from_string("10183A")
        strip.line.fill.background()
        strip.shadow.inherit = False
        # Split "FinOps Toolkit + Azure SRE Agent product team" by " + "
        parts = re.split(r"\s*\+\s*", cob, maxsplit=1)
        left_text = parts[0].strip() if parts else ""
        right_text = parts[1].strip() if len(parts) > 1 else ""

        # Left mark
        ltb = slide.shapes.add_textbox(Inches(0.7), strip_top, sl_w // 2 - Inches(0.7), strip_h)
        ltf = ltb.text_frame
        ltf.margin_left = ltf.margin_right = ltf.margin_top = ltf.margin_bottom = 0
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.LEFT
        lp_lbl = lp.add_run()
        lp_lbl.text = LABELS["cover.built-by-prefix"] + "  "
        lp_lbl.font.size = Pt(10)
        lp_lbl.font.color.rgb = RGBColor.from_string("8FA0CC")
        lp_run = lp.add_run()
        lp_run.text = left_text
        lp_run.font.size = Pt(15)
        lp_run.font.bold = True
        lp_run.font.color.rgb = RGBColor.from_string("FFFFFF")
        # Vertical center
        from pptx.enum.text import MSO_ANCHOR
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Right mark — only if there's a real second attribution after " + "
        if right_text:
            # If right_text is a URL (aka.ms/..., http://, https://), drop the
            # "In partnership with" prefix entirely — the URL stands on its own.
            is_url = bool(re.match(r"^(https?://|aka\.ms/|www\.)", right_text, re.IGNORECASE))
            rtb = slide.shapes.add_textbox(sl_w // 2, strip_top, sl_w // 2 - Inches(0.7), strip_h)
            rtf = rtb.text_frame
            rtf.margin_left = rtf.margin_right = rtf.margin_top = rtf.margin_bottom = 0
            rp = rtf.paragraphs[0]
            rp.alignment = PP_ALIGN.RIGHT
            if not is_url:
                rp_lbl = rp.add_run()
                rp_lbl.text = LABELS["cover.partnership-prefix"] + "  "
                rp_lbl.font.size = Pt(10)
                rp_lbl.font.color.rgb = RGBColor.from_string("8FA0CC")
            rp_run = rp.add_run()
            rp_run.text = right_text
            rp_run.font.size = Pt(15)
            rp_run.font.bold = True
            rp_run.font.color.rgb = RGBColor.from_string("FFFFFF")
            rtf.vertical_anchor = MSO_ANCHOR.MIDDLE


def render_bullets(slide, row):
    if slide.shapes.title is not None:
        slide.shapes.title.text = row["title"]
        # If title is long, shrink the title font to keep it on one line
        if len(row["title"]) > 50:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)
    body = get_placeholder(slide, 10)
    if body is not None:
        body.left = Inches(0.64)
        body.top = Inches(1.25)
        body.width = Inches(12.05)
        body.height = Inches(5.6)
    paras = [_format_bullet_with_tools(b, size=20) for b in row["bullets"]]
    _set_text_frame_with_runs(body, paras)


# ─── Footer chip: ask-coverage pill at slide bottom ──────────

def render_footer_chip(slide, row):
    """Render the ask-coverage footer chip at the bottom of a body slide."""
    addresses = row.get("addresses", [])
    verdict = row.get("verdict")
    if not addresses or not verdict:
        return

    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    sl_w = slide.part.package.presentation_part.presentation.slide_width
    sl_h = slide.part.package.presentation_part.presentation.slide_height

    colors = VERDICT_COLORS.get(verdict, VERDICT_COLORS["green"])
    text_hex, bg_hex = colors

    ask_text = LABELS["ask.footer-prefix"] + " " + LABELS["ask.footer-first-mark"] + LABELS["ask.footer-separator"].join(str(a) for a in addresses)
    verdict_emoji = {"green": "🟢", "yellow": "🟡", "red": "🔴"}.get(verdict, "")
    chip_text = f"{verdict_emoji} {ask_text}"

    chip_w = Inches(3.5)
    chip_h = Inches(0.32)
    chip_x = sl_w - chip_w - Inches(0.4)
    chip_y = sl_h - chip_h - Inches(0.15)

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, chip_x, chip_y, chip_w, chip_h)
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor.from_string(bg_hex)
    pill.line.color.rgb = RGBColor.from_string(text_hex)
    pill.line.width = Pt(0.75)
    pill.shadow.inherit = False

    tb = slide.shapes.add_textbox(chip_x, chip_y, chip_w, chip_h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = chip_text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(text_hex)


def render_cluster_chyron(slide, row, top=Inches(1.25)):
    """Render the cluster name as an uppercase blue chyron above body content.

    Matches the editorial signature established by render_ask_a:
    - Cluster name in blue uppercase, small bold
    - Used on ASK_A, ASK_B, ASK_C for visual consistency across the cluster's 3 slides
    """
    cluster = row.get("stage", "").strip()
    if not cluster:
        return

    sl_w = slide.part.package.presentation_part.presentation.slide_width
    margin_x = Inches(0.7)

    label_tb = slide.shapes.add_textbox(margin_x, top, sl_w - 2 * margin_x, Inches(0.35))
    ltf = label_tb.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_right = 0
    ltf.margin_top = ltf.margin_bottom = 0
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = cluster.upper()
    lr.font.size = Pt(11)
    lr.font.bold = True
    lr.font.color.rgb = RGBColor.from_string("0F6CBD")


def render_accent_rule(slide, top=Inches(1.85), height=Inches(5.6)):
    """Render the vertical blue accent rule on the left edge of the body content."""
    from pptx.enum.shapes import MSO_SHAPE
    margin_x = Inches(0.7)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   margin_x, top,
                                   Inches(0.06), height)
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string("0F6CBD")
    rule.line.fill.background()
    rule.shadow.inherit = False


# ─── ASK_A renderer: editorial pull-quote treatment ──────────

def render_ask_a(slide, row):
    """ASK slide: professional MCAPS pull-quote design.

    Title at top. Large italic quote(s) centered vertically with a left blue
    accent rule. Cluster name as a small label above the quote (chyron style).
    Footer chip stays. Removes default body placeholder; uses custom textboxes
    for full layout control.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    if slide.shapes.title is not None:
        slide.shapes.title.text = row["title"]
        if len(row["title"]) > 50:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)

    # Hide the default body placeholder — we draw our own pull-quote layout
    body_ph = get_placeholder(slide, 10)
    if body_ph is not None:
        sp = body_ph._element
        sp.getparent().remove(sp)

    # Separate quotes from attribution
    quotes = []
    attrib = []
    for b in row["bullets"]:
        if b.startswith(">") or b.startswith('"') or b.startswith("\u201c"):
            clean = b.lstrip("> ").strip('" \u201c\u201d')
            quotes.append(clean)
        else:
            attrib.append(b)
    if not quotes:
        quotes = row["bullets"][:2]
        attrib = row["bullets"][2:]

    # Limit to top 2 quotes for visual focus
    quotes = quotes[:2]

    # Adaptive font size based on combined quote length
    total_chars = sum(len(q) for q in quotes)
    if total_chars < 80:
        qsize = 36
    elif total_chars < 160:
        qsize = 30
    elif total_chars < 280:
        qsize = 24
    else:
        qsize = 20

    sl_w = slide.part.package.presentation_part.presentation.slide_width
    sl_h = slide.part.package.presentation_part.presentation.slide_height

    # Layout geometry: pull-quote centered, left accent rule
    margin_x = Inches(1.2)
    quote_left = margin_x + Inches(0.5)
    quote_w = sl_w - quote_left - Inches(1.2)

    # Quote occupies the middle band — vertically centered
    quote_top = Inches(2.2)
    quote_h = Inches(3.5)

    # Left accent rule (vertical blue bar)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   margin_x, quote_top + Inches(0.2),
                                   Inches(0.06), quote_h - Inches(0.4))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string("0F6CBD")
    rule.line.fill.background()
    rule.shadow.inherit = False

    # Cluster name label — small chyron above the quote (uppercase, light)
    cluster = row.get("stage", "").strip()
    if cluster:
        label_tb = slide.shapes.add_textbox(quote_left, Inches(1.7), quote_w, Inches(0.35))
        ltf = label_tb.text_frame
        ltf.word_wrap = True
        ltf.margin_left = ltf.margin_right = 0
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = cluster.upper()
        lr.font.size = Pt(11)
        lr.font.bold = True
        lr.font.color.rgb = RGBColor.from_string("0F6CBD")

    # Quote text — italic, large, vertically centered in the band
    qtb = slide.shapes.add_textbox(quote_left, quote_top, quote_w, quote_h)
    qtf = qtb.text_frame
    qtf.word_wrap = True
    qtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    qtf.margin_left = qtf.margin_right = Inches(0.1)
    qtf.margin_top = qtf.margin_bottom = 0

    for i, q in enumerate(quotes):
        p = qtf.paragraphs[0] if i == 0 else qtf.add_paragraph()
        p.space_after = Pt(qsize // 2)
        run = p.add_run()
        run.text = f"\u201c{q}\u201d"
        run.font.size = Pt(qsize)
        run.font.italic = True
        run.font.color.rgb = RGBColor.from_string("1B1B1F")

    # Attribution line below the quote (small, muted)
    if attrib:
        atb = slide.shapes.add_textbox(quote_left, quote_top + quote_h + Inches(0.1),
                                        quote_w, Inches(0.4))
        atf = atb.text_frame
        atf.word_wrap = True
        ap = atf.paragraphs[0]
        ar = ap.add_run()
        ar.text = f"— {attrib[0]}"
        ar.font.size = Pt(13)
        ar.font.color.rgb = RGBColor.from_string("595959")

    render_footer_chip(slide, row)


# ─── ASK_B renderer: editorial answer treatment ──────────────

def render_ask_b(slide, row):
    """SHOW slide: matches the editorial design language.

    Cluster chyron at top, blue accent rule on left, answer text (left ~40%)
    with bolded tool names, chart panel (right ~55%) with subtle border.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    if slide.shapes.title is not None:
        slide.shapes.title.text = row["title"]
        if len(row["title"]) > 50:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)

    # Remove default body placeholder — we draw a custom two-column layout
    body_ph = get_placeholder(slide, 10)
    if body_ph is not None:
        sp = body_ph._element
        sp.getparent().remove(sp)

    # Cluster chyron + accent rule (matches ASK_A and ASK_C signature)
    render_cluster_chyron(slide, row, top=Inches(1.25))
    render_accent_rule(slide, top=Inches(1.85), height=Inches(5.4))

    sl_w = slide.part.package.presentation_part.presentation.slide_width
    sl_h = slide.part.package.presentation_part.presentation.slide_height

    # Layout geometry: text left (~38%), chart right (~58%)
    content_top = Inches(1.85)
    content_h = Inches(4.9)

    text_left = Inches(1.0)
    text_w = Inches(4.6)

    chart_left = text_left + text_w + Inches(0.3)
    chart_w = sl_w - chart_left - Inches(0.7)

    # Determine layout strategy: chart, owner-card, or full-width text-only
    assets = row.get("assets", {})
    img_file = assets.get("image")
    has_chart = False

    # Pre-check: if no chart asset, decide whether right panel makes sense
    owner_pattern = re.compile(r"(owns|owner|stays with|out of scope|not in this release|this release does not)", re.IGNORECASE)
    owner_bullets = [b for b in row["bullets"] if owner_pattern.search(b)] if not img_file else []
    needs_two_col = bool(img_file) or bool(owner_bullets)

    if needs_two_col:
        # Two-column layout: left text + right chart/card
        ttb = slide.shapes.add_textbox(text_left, content_top, text_w, content_h)
        ttf = ttb.text_frame
        ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = 0
        ttf.margin_top = Inches(0.1)

        bullets = row["bullets"][:5]
        for i, b in enumerate(bullets):
            p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
            p.space_after = Pt(8)
            para_data = _format_bullet_with_tools(b, size=14)
            for r in list(p.runs):
                r.text = ""
            for run_data in para_data["runs"]:
                run = p.add_run()
                run.text = run_data.get("text", "") or ""
                font = run.font
                if run_data.get("size"):
                    font.size = Pt(run_data["size"])
                if run_data.get("bold"):
                    font.bold = True
                if run_data.get("color"):
                    font.color.rgb = RGBColor.from_string(run_data["color"])
    else:
        # Full-width text — no left/right split
        full_w = sl_w - text_left - Inches(0.7)
        ttb = slide.shapes.add_textbox(text_left, content_top, full_w, content_h)
        ttf = ttb.text_frame
        ttf.word_wrap = True
        ttf.margin_left = ttf.margin_right = 0
        ttf.margin_top = Inches(0.1)
        for i, b in enumerate(row["bullets"][:8]):
            p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
            p.space_after = Pt(8)
            para_data = _format_bullet_with_tools(b, size=14)
            for r in list(p.runs):
                r.text = ""
            for run_data in para_data["runs"]:
                run = p.add_run()
                run.text = run_data.get("text", "") or ""
                font = run.font
                if run_data.get("size"):
                    font.size = Pt(run_data["size"])
                if run_data.get("bold"):
                    font.bold = True
                if run_data.get("color"):
                    font.color.rgb = RGBColor.from_string(run_data["color"])

    # Chart panel: subtle border around the image
    assets = row.get("assets", {})
    img_file = assets.get("image")
    has_chart = False

    if img_file:
        img_path = CHARTS / img_file
        if not img_path.exists():
            img_path = ASSETS / img_file
        if img_path.exists():
            if img_path.suffix.lower() == ".svg":
                png_path = img_path.with_suffix(".rasterized.png")
                if not png_path.exists():
                    import subprocess
                    subprocess.run(["rsvg-convert", "-w", "1600", "-o", str(png_path), str(img_path)], check=True)
                img_path = png_path

            # Read actual image dimensions and compute fitted size with aspect preservation
            try:
                from PIL import Image as PILImage
                with PILImage.open(str(img_path)) as pil_img:
                    img_w_px, img_h_px = pil_img.size
                img_aspect = img_w_px / img_h_px
                panel_aspect = chart_w / content_h

                if img_aspect > panel_aspect:
                    # Image is wider — fit to width
                    fit_w = chart_w
                    fit_h = int(chart_w / img_aspect)
                    fit_left = chart_left
                    fit_top = content_top + (content_h - fit_h) // 2
                else:
                    # Image is taller — fit to height
                    fit_h = content_h
                    fit_w = int(content_h * img_aspect)
                    fit_top = content_top
                    fit_left = chart_left + (chart_w - fit_w) // 2

                # Tight frame matching the actual fitted image (Tier 1B fix)
                # This eliminates the oversized empty panel around small charts
                frame_pad = Inches(0.05)
                frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                fit_left - frame_pad,
                                                fit_top - frame_pad,
                                                fit_w + 2 * frame_pad,
                                                fit_h + 2 * frame_pad)
                frame.fill.solid()
                frame.fill.fore_color.rgb = RGBColor.from_string("FAFAFA")
                frame.line.color.rgb = RGBColor.from_string("E5E5EA")
                frame.line.width = Pt(0.5)
                frame.shadow.inherit = False

                slide.shapes.add_picture(str(img_path), fit_left, fit_top,
                                         width=fit_w, height=fit_h)
                has_chart = True
            except Exception as e:
                print(f"   ! ASK_B image embed failed: {e}")
        else:
            print(f"   ! ASK_B chart missing: {img_file}")

    # If no chart but we have owner-attribution bullets, render them as a card on the right
    if not has_chart and owner_bullets:
        display_bullets = owner_bullets[:6]
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       chart_left - Inches(0.1),
                                       content_top - Inches(0.05),
                                       chart_w + Inches(0.2),
                                       content_h + Inches(0.1))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("F4F6FB")
        card.line.color.rgb = RGBColor.from_string("D2D2D7")
        card.line.width = Pt(0.5)
        card.shadow.inherit = False

        # Label
        ltb = slide.shapes.add_textbox(chart_left + Inches(0.2),
                                        content_top + Inches(0.2),
                                        chart_w - Inches(0.4), Inches(0.4))
        ltf = ltb.text_frame
        ltf.margin_left = ltf.margin_right = 0
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = LABELS["ask-b.owners-header"].upper()
        lr.font.size = Pt(11)
        lr.font.bold = True
        lr.font.color.rgb = RGBColor.from_string("0F6CBD")

        otb = slide.shapes.add_textbox(chart_left + Inches(0.2),
                                        content_top + Inches(0.7),
                                        chart_w - Inches(0.4), content_h - Inches(0.9))
        otf = otb.text_frame
        otf.word_wrap = True
        for i, b in enumerate(display_bullets):
            p = otf.paragraphs[0] if i == 0 else otf.add_paragraph()
            p.space_after = Pt(8)
            para_data = _format_bullet_with_tools(b, size=13)
            for r in list(p.runs):
                r.text = ""
            for run_data in para_data["runs"]:
                run = p.add_run()
                run.text = run_data.get("text", "") or ""
                font = run.font
                if run_data.get("size"):
                    font.size = Pt(run_data["size"])
                if run_data.get("bold"):
                    font.bold = True
                if run_data.get("color"):
                    font.color.rgb = RGBColor.from_string(run_data["color"])

    render_footer_chip(slide, row)


def _format_bullet_with_tools(text, size=14):
    """Format a bullet by detecting tool/task names and bolding them.

    Tool/task pattern: lowercase identifiers with hyphens (e.g., 'capacity-daily-monitor',
    'deploy-budget', 'vm-quota-usage'). Returns a paragraph dict with multiple runs.
    """
    # Pattern matches identifiers with at least one hyphen, lowercase letters
    pattern = re.compile(r"\b([a-z][a-z0-9]*(?:-[a-z0-9]+){1,})\b")
    runs = []
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            runs.append({"text": text[pos:m.start()], "size": size, "color": "1B1B1F"})
        runs.append({"text": m.group(0), "size": size, "color": "0F6CBD", "bold": True})
        pos = m.end()
    if pos < len(text):
        runs.append({"text": text[pos:], "size": size, "color": "1B1B1F"})
    if not runs:
        runs = [{"text": text, "size": size, "color": "1B1B1F"}]
    return {"runs": runs, "space_after": 6}


def _set_text_frame_with_runs(ph, paragraphs):
    """Set text frame from paragraphs that may contain multiple styled runs."""
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, p in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        for r in list(para.runs):
            r.text = ""
        if p.get("space_after"):
            para.space_after = Pt(p["space_after"])
        runs_data = p.get("runs", [{"text": p.get("text", "")}])
        for run_data in runs_data:
            run = para.add_run()
            run.text = run_data.get("text", "") or ""
            font = run.font
            if run_data.get("size"):
                font.size = Pt(run_data["size"])
            if run_data.get("bold"):
                font.bold = True
            if run_data.get("color"):
                font.color.rgb = RGBColor.from_string(run_data["color"])


# ─── ASK_C renderer: editorial details treatment ─────────────

def render_ask_c(slide, row):
    """TELL slide: matches the editorial pull-quote design language.

    Cluster chyron at top (uppercase blue), blue accent rule on left,
    technical content in clean hierarchy:
      - Prose lines: serif body
      - Code lines (between ``` markers): dark monospace blocks
      - Monday move: highlighted blue callout panel at the bottom
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    if slide.shapes.title is not None:
        slide.shapes.title.text = row["title"]
        if len(row["title"]) > 50:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(24)

    # Remove default body placeholder — we draw a custom layout
    body_ph = get_placeholder(slide, 10)
    if body_ph is not None:
        sp = body_ph._element
        sp.getparent().remove(sp)

    # Cluster chyron at top (matches ASK_A signature)
    render_cluster_chyron(slide, row, top=Inches(1.25))
    # Note: accent rule is drawn AT END to match actual content height

    sl_w = slide.part.package.presentation_part.presentation.slide_width
    sl_h = slide.part.package.presentation_part.presentation.slide_height
    content_left = Inches(1.0)
    content_w = sl_w - content_left - Inches(0.7)

    # Group bullets into segments: prose, code (monospace block), monday
    segments = []
    in_code = False
    code_buf = []
    for b in row["bullets"]:
        if b.strip() == "```":
            if in_code:
                if code_buf:
                    segments.append({"type": "code", "lines": code_buf})
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buf.append(b)
            continue
        is_monday = LABELS["ask-c.monday-header"].lower() in b.lower() or b.strip().startswith(LABELS["ask-c.code-marker-prefix"])
        if is_monday:
            segments.append({"type": "monday", "text": b})
        else:
            segments.append({"type": "prose", "text": b})
    if in_code and code_buf:
        segments.append({"type": "code", "lines": code_buf})

    # Lay out segments top-to-bottom with proper spacing
    y = Inches(1.85)
    bottom_limit = sl_h - Inches(0.8)
    gap = Inches(0.18)

    for seg in segments:
        if y >= bottom_limit:
            break

        if seg["type"] == "code":
            # Dark code block — fixed-width monospace, preserves alignment
            n_lines = len(seg["lines"])
            line_h = Inches(0.22)
            block_h = line_h * n_lines + Inches(0.25)
            if y + block_h > bottom_limit:
                block_h = bottom_limit - y

            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         content_left, y, content_w, block_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor.from_string("1E1E2E")
            bg.line.color.rgb = RGBColor.from_string("3A3A4A")
            bg.line.width = Pt(0.5)
            bg.shadow.inherit = False

            tb = slide.shapes.add_textbox(content_left + Inches(0.2),
                                           y + Inches(0.12),
                                           content_w - Inches(0.4),
                                           block_h - Inches(0.24))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0

            for i, line in enumerate(seg["lines"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = line
                run.font.size = Pt(11)
                run.font.name = "Consolas"
                run.font.color.rgb = RGBColor.from_string("D4D4D4")

            y += block_h + gap

        elif seg["type"] == "monday":
            # Highlighted Monday-move callout — narrower, centered, sized to fit single line
            callout_h = Inches(0.45)
            if y + callout_h > bottom_limit:
                callout_h = bottom_limit - y

            # Make callout narrower than the full content width — feels less heavy
            callout_w = min(content_w, Inches(8.5))
            callout_x = content_left + (content_w - callout_w) // 2

            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         callout_x, y, callout_w, callout_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor.from_string("EAF3FB")
            bg.line.color.rgb = RGBColor.from_string("0F6CBD")
            bg.line.width = Pt(1.0)
            bg.shadow.inherit = False

            tb = slide.shapes.add_textbox(callout_x + Inches(0.25),
                                           y + Inches(0.05),
                                           callout_w - Inches(0.5),
                                           callout_h - Inches(0.1))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = seg["text"]
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor.from_string("0F6CBD")

            y += callout_h + gap

        else:  # prose
            # Estimate prose height: ~0.30" per ~80 chars
            n_visual_lines = max(1, (len(seg["text"]) // 90) + 1)
            prose_h = Inches(0.30) * n_visual_lines
            if y + prose_h > bottom_limit:
                prose_h = bottom_limit - y

            tb = slide.shapes.add_textbox(content_left, y, content_w, prose_h)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0

            # Format prose with tool-name bolding
            para_data = _format_bullet_with_tools(seg["text"], size=14)
            p = tf.paragraphs[0]
            for r in list(p.runs):
                r.text = ""
            for run_data in para_data["runs"]:
                run = p.add_run()
                run.text = run_data.get("text", "") or ""
                font = run.font
                if run_data.get("size"):
                    font.size = Pt(run_data["size"])
                if run_data.get("bold"):
                    font.bold = True
                if run_data.get("color"):
                    font.color.rgb = RGBColor.from_string(run_data["color"])

            y += prose_h + Inches(0.18)

    # Draw accent rule sized to actual content height (Tier 1A fix)
    actual_content_h = y - Inches(1.45)
    rule_h = max(Inches(1.0), min(actual_content_h, Inches(5.4)))
    render_accent_rule(slide, top=Inches(1.85), height=rule_h)

    render_footer_chip(slide, row)


def _set_text_frame_with_mono(ph, paragraphs):
    """Like set_text_frame but supports `mono: True` for monospace font."""
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, p in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        for r in list(para.runs):
            r.text = ""

        run = para.add_run()
        run.text = p.get("text", "") or ""

        font = run.font
        if p.get("size"):
            font.size = Pt(p["size"])
        if p.get("bold"):
            font.bold = True
        if p.get("mono"):
            font.name = "Consolas"
        if p.get("color"):
            font.color.rgb = RGBColor.from_string(p["color"])


# ─── INDEX slide renderer: ask map for the deck ──────────────

def render_index(slide, row):
    """Index slide: clean two-column list of clusters with ask numbers + verdicts.

    Reads bullets from the outline directly. Each bullet is parsed as:
      "<cluster_id> · <verdict> · <name> · <comma,asks>"
    e.g. "P1.1 · green · Quota ≠ Capacity · 8,13,14,23,24"
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    if slide.shapes.title is not None:
        slide.shapes.title.text = row["title"]

    # Remove default body placeholder — custom layout
    body_ph = get_placeholder(slide, 10)
    if body_ph is not None:
        sp = body_ph._element
        sp.getparent().remove(sp)

    sl_w = slide.part.package.presentation_part.presentation.slide_width
    sl_h = slide.part.package.presentation_part.presentation.slide_height
    margin_x = Inches(0.7)
    list_top = Inches(1.3)

    # Parse bullets into entries
    entries = []
    for b in row.get("bullets", []):
        parts = [p.strip() for p in b.split("·")]
        if len(parts) < 4:
            continue
        cluster_id, verdict, name, asks_str = parts[0], parts[1], parts[2], parts[3]
        asks = [int(a.strip()) for a in asks_str.split(",") if a.strip().isdigit()]
        entries.append({
            "cluster_id": cluster_id,
            "verdict": verdict.lower(),
            "name": name,
            "asks": asks,
        })

    n = len(entries)
    if n == 0:
        return

    # Two columns side-by-side, split entries roughly evenly
    col_w = (sl_w - 2 * margin_x - Inches(0.3)) // 2
    col_left_x = margin_x
    col_right_x = margin_x + col_w + Inches(0.3)

    half = (n + 1) // 2
    left_entries = entries[:half]
    right_entries = entries[half:]

    list_h = sl_h - list_top - Inches(0.8)
    row_h = list_h // max(half, 1)

    verdict_dot = {"green": "🟢", "yellow": "🟡", "red": "🔴"}

    def draw_entry(x, y, w, h, entry):
        v = entry.get("verdict", "green")
        vdot = verdict_dot.get(v, "")
        cid = entry.get("cluster_id", "")
        name = entry.get("name", "")
        asks = entry.get("asks", [])
        ask_str = LABELS["index.entry-separator"].join(LABELS["index.ask-mark"] + str(a) for a in asks)

        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = Inches(0.04)

        p1 = tf.paragraphs[0]
        p1.space_after = Pt(0)
        r_dot = p1.add_run()
        r_dot.text = f"{vdot}  "
        r_dot.font.size = Pt(11)

        r_cid = p1.add_run()
        r_cid.text = f"{cid}  "
        r_cid.font.size = Pt(11)
        r_cid.font.bold = True
        r_cid.font.color.rgb = RGBColor.from_string("0F6CBD")

        r_name = p1.add_run()
        r_name.text = name
        r_name.font.size = Pt(13)
        r_name.font.color.rgb = RGBColor.from_string("1B1B1F")

        p2 = tf.add_paragraph()
        p2.space_before = Pt(0)
        p2.space_after = Pt(0)
        r_asks = p2.add_run()
        r_asks.text = f"   Asks {ask_str}"
        r_asks.font.size = Pt(10)
        r_asks.font.color.rgb = RGBColor.from_string("595959")

    for i, e in enumerate(left_entries):
        draw_entry(col_left_x, list_top + i * row_h, col_w, row_h, e)
    for i, e in enumerate(right_entries):
        draw_entry(col_right_x, list_top + i * row_h, col_w, row_h, e)


RENDERERS = {
    "TITLE": render_title,
    "BULLETS": render_bullets,
    "ASK_A": render_ask_a,
    "ASK_B": render_ask_b,
    "ASK_C": render_ask_c,
    "INDEX": render_index,
}


# ─── Slide deletion ──────────────────────────────────────────

def delete_all_slides(prs):
    """Delete every slide from the presentation cleanly."""
    sldIdLst = prs.slides._sldIdLst
    rIds = [sldId.rId for sldId in list(sldIdLst)]
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    # Drop the relationships and parts
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ─── Main pipeline ───────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Build V8 deck pptx from canonical outline.")
    parser.add_argument(
        "--part",
        choices=["1", "2", "3", "all"],
        default="all",
        help="Build only the slides for one part (1=deal motion, 2=operate motion, 3=honest+close). Default 'all' produces the unified deck.",
    )
    args = parser.parse_args()

    rows = parse_v8()
    print(f"Parsed {len(rows)} V8 rows\n")

    if args.part != "all":
        pattern, out_name = PART_FILTERS[args.part]
        before = len(rows)
        rows = [r for r in rows if pattern.match(r["num"])]
        print(f"Filtered to part {args.part}: {len(rows)} of {before} slides match {pattern.pattern}")
        out_path = ROOT / out_name
    else:
        out_path = OUT

    prs = Presentation(str(SOURCE))

    # Validate layouts exist before doing anything destructive
    for kind in {r["kind"] for r in rows}:
        find_layout(prs, KIND_LAYOUT[kind])

    delete_all_slides(prs)
    print(f"Cleared all original slides.\n")

    for row in rows:
        kind = row["kind"]
        layout = find_layout(prs, KIND_LAYOUT[kind])
        slide = prs.slides.add_slide(layout)

        # Render content
        RENDERERS[kind](slide, row)

        # Speaker notes
        if row["notes"] and any(n.strip() for n in row["notes"]):
            tf = slide.notes_slide.notes_text_frame
            tf.clear()
            tf.paragraphs[0].text = row["notes"][0]
            for note_para in row["notes"][1:]:
                p = tf.add_paragraph()
                p.text = note_para

        print(f"  {row['num']:5s}  [{kind:18s}]  {layout.name:40s}  {row['title'][:55]}")

    # Scrub author metadata via core_props
    prs.core_properties.author = LABELS["pptx.author"]
    prs.core_properties.last_modified_by = LABELS["pptx.author"]

    # Safe save: detect PowerPoint lock file and use atomic write to avoid
    # corrupting the file mid-read for any process holding it open.
    out_dir = out_path.parent
    out_name = out_path.name
    lock_file = out_dir / f"~${out_name}"

    if lock_file.exists():
        # PowerPoint has the file open. Writing to out_path will desync its view
        # and cause a "needs repair" prompt next time the user clicks anything.
        # Refuse to overwrite — write to a sibling and tell the user.
        import time
        alt_name = out_path.stem + f".rebuild-{int(time.time())}.pptx"
        alt_out = out_dir / alt_name
        prs.save(str(alt_out))
        print(f"\n⚠️  PowerPoint lock file detected ({lock_file.name})")
        print(f"   {out_path.name} is currently open in PowerPoint.")
        print(f"   To avoid corrupting your open copy, wrote rebuild to:")
        print(f"   {alt_out}")
        print(f"   Close PowerPoint and re-run, or open the new file directly.")
        print(f"Slides: {len(prs.slides)}")
        return

    # Atomic write: temp file → fsync → rename. This guarantees readers
    # never see a partially-written file.
    import os, tempfile
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx", dir=str(out_dir))
    os.close(tmp_fd)
    try:
        prs.save(tmp_path)
        # Fsync to ensure bytes are on disk before rename
        with open(tmp_path, "rb") as f:
            os.fsync(f.fileno())
        os.replace(tmp_path, str(out_path))
        print(f"\nWrote {out_path}")
        print(f"Slides: {len(prs.slides)}")
    except Exception:
        # Clean up the temp file on failure
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise


if __name__ == "__main__":
    main()
