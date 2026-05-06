"""Build the V8 deck from per-slide YAML files in slides/.

Mirrors build.py's pipeline exactly except the canonical source: instead of
parsing deck-outline-v8.md, this reads slides/*.yaml in numeric prefix order,
translates each into the row dict that build.py's renderers expect, then
delegates to build.py's RENDERERS dispatch table.

Usage:
    python build_yaml.py             # unified deck (94 slides)
    python build_yaml.py --part 1    # part 1 only
    python build_yaml.py --part 2    # part 2 only
    python build_yaml.py --part 3    # part 3 only
    python build_yaml.py --output foo.pptx
"""
import argparse
import os
import re
import sys
import tempfile
import time
from pathlib import Path

import yaml
from pptx import Presentation

import build as _b

ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / "slides"


def parse_slides():
    """Read slides/*.yaml in filename order, return rows in build.py schema."""
    rows = []
    files = sorted(SLIDES_DIR.glob("*.yaml"))
    if not files:
        sys.exit(f"No YAML files found in {SLIDES_DIR}/")
    for f in files:
        s = yaml.safe_load(f.read_text(encoding="utf-8"))
        bullets = [_b.strip_md(_b.strip_bullet(b)) for b in _b.html_br_to_lines(s.get("content", "") or "")]
        notes_text = s.get("notes", "") or ""
        notes_paras = [_b.strip_md(p) for p in _b.html_br_to_lines(notes_text)] or [_b.strip_md(notes_text)]
        layout = s.get("layout", {}) or {}
        kind = layout.get("kind", "BULLETS")
        addresses = layout.get("addresses") or []
        verdict = layout.get("verdict")
        # Assets from screens directive
        assets = {}
        screens = s.get("screens") or ""
        if screens:
            am = re.search(r"assets:\s*(.+?)$", screens)
            if am:
                for kv in re.finditer(r"(\w+)\s*=\s*([^\s|]+)", am.group(1)):
                    assets[kv.group(1)] = kv.group(2)
        rows.append({
            "num": s["id"],
            "stage": _b.strip_md(s.get("cluster", "") or ""),
            "title": _b.strip_md(s.get("title", "") or ""),
            "bullets": bullets,
            "notes": notes_paras,
            "kind": kind,
            "layout": "",
            "assets": assets,
            "addresses": addresses,
            "verdict": verdict.lower() if verdict else None,
        })
    return rows


def main():
    parser = argparse.ArgumentParser(description="Build V8 deck pptx from per-slide YAML.")
    parser.add_argument("--part", choices=["1", "2", "3", "all"], default="all")
    parser.add_argument("--output", help="Override output filename")
    args = parser.parse_args()

    rows = parse_slides()

    if args.part != "all":
        pattern, out_name = _b.PART_FILTERS[args.part]
        prog = re.compile(pattern)
        before = len(rows)
        rows = [r for r in rows if prog.match(r["num"])]
        print(f"Filtered to part {args.part}: {len(rows)} of {before} slides match {prog.pattern}")
        out_path = ROOT / out_name
    else:
        out_path = _b.OUT

    if args.output:
        out_path = Path(args.output) if Path(args.output).is_absolute() else ROOT / args.output

    if not rows:
        sys.exit(f"No rows matched filter for part {args.part}")

    prs = Presentation(str(_b.SOURCE))

    # Validate every kind we will render has a known layout
    for kind in {r["kind"] for r in rows}:
        _b.find_layout(prs, _b.KIND_LAYOUT[kind])

    _b.delete_all_slides(prs)
    print(f"Cleared all original slides.\n")

    for row in rows:
        kind = row["kind"]
        layout = _b.find_layout(prs, _b.KIND_LAYOUT[kind])
        slide = prs.slides.add_slide(layout)
        _b.RENDERERS[kind](slide, row)

        if row["notes"] and any(n.strip() for n in row["notes"]):
            tf = slide.notes_slide.notes_text_frame
            tf.clear()
            tf.paragraphs[0].text = row["notes"][0]
            for note_para in row["notes"][1:]:
                p = tf.add_paragraph()
                p.text = note_para

        print(f"  {row['num']:5s}  [{kind:18s}]  {layout.name:40s}  {row['title'][:55]}")

    # Scrub author metadata
    prs.core_properties.author = _b.LABELS["pptx.author"]
    prs.core_properties.last_modified_by = _b.LABELS["pptx.author"]

    # Lock-file detection
    out_dir = out_path.parent
    out_name_final = out_path.name
    lock = out_dir / f"~${out_name_final}"
    if lock.exists():
        alt = out_path.stem + f".rebuild-{int(time.time())}.pptx"
        out_path = out_dir / alt
        print(f"\n⚠️  PowerPoint lock detected on {out_name_final}. Writing to {alt} instead.")

    # Atomic save
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx", dir=str(out_dir))
    os.close(tmp_fd)
    try:
        prs.save(tmp_path)
        os.replace(tmp_path, str(out_path))
        print(f"\nWrote {out_path}")
        print(f"Slides: {len(rows)}")
    except Exception:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise


if __name__ == "__main__":
    main()
