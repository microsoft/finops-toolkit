#!/usr/bin/env python3
"""Pixel-perfect rubric checker for the V8 release training deck.

Hard rules — every check is binary. A slide either passes or fails.
No vibes, no judgement calls. The script catches what humans miss.

Usage:
    python3 rubric.py                    # check all 94 slides
    python3 rubric.py --verbose          # show every shape on every slide
    python3 rubric.py --slide N          # check just slide N

Exit code is 0 if all rules pass; 1 if any rule fails.
"""
from __future__ import annotations

import sys
from pathlib import Path
from typing import NamedTuple

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
DECK = ROOT / "finops-toolkit-sre-agent-release-training.pptx"

# Canvas dimensions (16:9 @ standard widescreen)
CANVAS_W = Emu(12192000)   # 13.33"
CANVAS_H = Emu(6858000)    #  7.50"

# Tolerance for "touching" — anything closer than this is a violation
MIN_GAP_EMU = Emu(91440)   # 0.10"


class Bbox(NamedTuple):
    """Axis-aligned bounding box in EMU."""
    left: int
    top: int
    right: int
    bottom: int
    name: str

    @classmethod
    def from_shape(cls, shape) -> "Bbox":
        try:
            l = shape.left or 0
            t = shape.top or 0
            w = shape.width or 0
            h = shape.height or 0
        except (AttributeError, TypeError):
            return cls(0, 0, 0, 0, getattr(shape, "name", "?"))
        return cls(l, t, l + w, t + h, getattr(shape, "name", "?"))

    def overlaps(self, other: "Bbox") -> bool:
        """True if the two bboxes intersect (share any pixel)."""
        if self.right <= other.left or other.right <= self.left:
            return False
        if self.bottom <= other.top or other.bottom <= self.top:
            return False
        return True

    def too_close(self, other: "Bbox", min_gap: int = MIN_GAP_EMU) -> bool:
        """True if the two bboxes are closer than min_gap and don't overlap."""
        if self.overlaps(other):
            return False
        h_gap = max(0, max(other.left - self.right, self.left - other.right))
        v_gap = max(0, max(other.top - self.bottom, self.top - other.bottom))
        # Only flag if both axes have content overlap (so they're stacked)
        h_overlap = not (self.right <= other.left or other.right <= self.left)
        v_overlap = not (self.bottom <= other.top or other.bottom <= self.top)
        if h_overlap and 0 < v_gap < min_gap:
            return True
        if v_overlap and 0 < h_gap < min_gap:
            return True
        return False

    def off_canvas(self) -> bool:
        """True if any part of the bbox extends outside the slide canvas."""
        return (
            self.left < 0
            or self.top < 0
            or self.right > CANVAS_W
            or self.bottom > CANVAS_H
        )

    def __repr__(self) -> str:
        return (
            f"<{self.name} l={self.left/914400:.2f}\" t={self.top/914400:.2f}\" "
            f"r={self.right/914400:.2f}\" b={self.bottom/914400:.2f}\">"
        )


# ─── Rubric definitions ──────────────────────────────────────

def rule_no_overlapping_text(slide) -> list[str]:
    """RULE 1: No two text-bearing shapes may overlap.

    Text shapes = placeholders + textboxes that have non-empty text frames.
    Pictures, lines, and decorative shapes are allowed to overlap with anything.
    """
    text_bboxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        text_bboxes.append((shape, Bbox.from_shape(shape), text[:60]))

    failures = []
    for i, (sa, ba, ta) in enumerate(text_bboxes):
        for sb, bb, tb in text_bboxes[i + 1:]:
            if ba.overlaps(bb):
                failures.append(
                    f"  OVERLAP: {ba.name!r} ({ta!r}) overlaps {bb.name!r} ({tb!r})"
                )
    return failures


def rule_no_text_off_canvas(slide) -> list[str]:
    """RULE 2: No text shape may extend off the slide canvas."""
    failures = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        bbox = Bbox.from_shape(shape)
        if bbox.off_canvas():
            failures.append(
                f"  OFF-CANVAS: {bbox.name!r} extends past 13.33\"x7.50\" canvas: {bbox}"
            )
    return failures


def rule_no_text_touching_text(slide) -> list[str]:
    """RULE 3: Text shapes must have at least 0.10\" gap between them on shared axis."""
    text_bboxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        text_bboxes.append((shape, Bbox.from_shape(shape)))

    failures = []
    for i, (sa, ba) in enumerate(text_bboxes):
        for sb, bb in text_bboxes[i + 1:]:
            if ba.too_close(bb):
                failures.append(
                    f"  TOUCHING: {ba.name!r} too close to {bb.name!r} (<0.10\" gap)"
                )
    return failures


def rule_no_circles_as_ovals(slide) -> list[str]:
    """RULE 4: Shapes labeled as circles/ovals must be square (1:1 aspect)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    failures = []
    for shape in slide.shapes:
        # MSO_SHAPE_TYPE.AUTO_SHAPE with auto_shape_type == OVAL
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if not hasattr(shape, "auto_shape_type"):
                continue
            from pptx.enum.shapes import MSO_SHAPE
            if shape.auto_shape_type != MSO_SHAPE.OVAL:
                continue
        except (AttributeError, ValueError):
            continue

        w = shape.width or 0
        h = shape.height or 0
        if w == 0 or h == 0:
            continue
        ratio = w / h
        if not (0.95 <= ratio <= 1.05):
            failures.append(
                f"  OVAL-NOT-CIRCLE: {shape.name!r} is {w/914400:.2f}\"x{h/914400:.2f}\" "
                f"(ratio {ratio:.2f}, must be 0.95-1.05)"
            )
    return failures


def rule_no_image_aspect_distortion(slide) -> list[str]:
    """RULE 5: Embedded images must preserve their native aspect ratio.

    Reads the actual image bytes and compares the embedded width:height
    against the source dimensions.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io

    failures = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            blob = shape.image.blob
            with Image.open(io.BytesIO(blob)) as img:
                src_w, src_h = img.size
        except Exception:
            continue
        embed_w = shape.width or 0
        embed_h = shape.height or 0
        if embed_w == 0 or embed_h == 0 or src_w == 0 or src_h == 0:
            continue
        src_ratio = src_w / src_h
        embed_ratio = embed_w / embed_h
        deviation = abs(src_ratio - embed_ratio) / src_ratio
        if deviation > 0.02:  # >2% distortion
            failures.append(
                f"  IMAGE-STRETCHED: {shape.name!r} native {src_ratio:.2f} "
                f"vs embedded {embed_ratio:.2f} ({deviation*100:.1f}% distortion)"
            )
    return failures


# Register rules in checking order
RULES = [
    ("R1: no overlapping text",       rule_no_overlapping_text),
    ("R2: no off-canvas text",        rule_no_text_off_canvas),
    ("R3: no touching text (<0.1\")", rule_no_text_touching_text),
    ("R4: ovals are circles",         rule_no_circles_as_ovals),
    ("R5: images preserve aspect",    rule_no_image_aspect_distortion),
]


def check_slide(slide, idx: int, verbose: bool = False) -> list[tuple[str, list[str]]]:
    """Run every rule against one slide. Returns list of (rule_name, failures)."""
    results = []
    for rule_name, rule_fn in RULES:
        try:
            failures = rule_fn(slide)
        except Exception as e:
            failures = [f"  RULE-ERROR: {type(e).__name__}: {e}"]
        results.append((rule_name, failures))

    if verbose:
        print(f"\n--- slide {idx + 1} shapes ---")
        for shape in slide.shapes:
            bbox = Bbox.from_shape(shape)
            text = ""
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = repr(shape.text_frame.text.strip()[:50])
            print(f"  {bbox} {text}")

    return results


def main() -> int:
    import argparse
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--verbose", action="store_true", help="dump every shape on every slide")
    ap.add_argument("--slide", type=int, help="only check one slide (1-indexed)")
    ap.add_argument("--deck", type=Path, default=DECK, help="path to .pptx")
    args = ap.parse_args()

    if not args.deck.exists():
        print(f"❌ Deck not found: {args.deck}")
        return 1

    prs = Presentation(str(args.deck))
    total_slides = len(prs.slides)
    target_slides = (
        [args.slide - 1]
        if args.slide is not None
        else range(total_slides)
    )

    print(f"=== Pixel-perfect rubric: {args.deck.name} ({total_slides} slides) ===\n")
    print("Rules:")
    for name, _ in RULES:
        print(f"  • {name}")
    print()

    total_failures = 0
    failed_slides = []
    for idx in target_slides:
        slide = list(prs.slides)[idx]
        results = check_slide(slide, idx, verbose=args.verbose)
        slide_failures = sum(len(f) for _, f in results)
        if slide_failures > 0:
            total_failures += slide_failures
            failed_slides.append(idx + 1)
            title = ""
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()[:50]
            print(f"❌ slide-{idx+1:02d} ({title}) — {slide_failures} failure(s)")
            for rule_name, failures in results:
                if failures:
                    print(f"   [{rule_name}]")
                    for f in failures:
                        print(f"   {f}")
        elif args.verbose:
            print(f"✅ slide-{idx+1:02d}")

    print()
    print(f"=== Summary ===")
    print(f"Slides checked:  {len(list(target_slides))}")
    print(f"Slides failing:  {len(failed_slides)}")
    print(f"Total failures:  {total_failures}")
    if failed_slides:
        print(f"Failed slides:   {', '.join(str(s) for s in failed_slides)}")
        print("\n❌ DECK FAILS RUBRIC")
        return 1
    print("\n✅ DECK PASSES RUBRIC")
    return 0


if __name__ == "__main__":
    sys.exit(main())
